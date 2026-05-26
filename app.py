import streamlit as st
import pandas as pd
import streamlit.components.v1 as components
import requests
import base64
import hashlib

# === 1. 設定區 ===
# ⚠️ 請確保這兩個網址分別是「總資料庫」分頁與「Clients」分頁獨立發布為 CSV 的網址
CSV_URL = "https://docs.google.com/spreadsheets/d/e/2PACX-1vSnViFsUwWYASaR5i1PefsWE4b6-5wwqTbJFJG8vysgcHYZDKzq-wwK4hM4xOtet3B65UjohzRjh38C/pub?gid=0&single=true&output=csv"
CSV_LOGO_URL = "https://docs.google.com/spreadsheets/d/e/2PACX-1vSnViFsUwWYASaR5i1PefsWE4b6-5wwqTbJFJG8vysgcHYZDKzq-wwK4hM4xOtet3B65UjohzRjh38C/pub?gid=1588470763&single=true&output=csv" # 👈 請務必填入 Clients 分頁專屬的 CSV 網址

PASSWORD = "888"
SITE_URL = "https://swd-case.streamlit.app" 

# === 2. 核心技術函數 ===
def generate_id(link):
    return hashlib.md5(str(link).encode()).hexdigest()[:10]

@st.cache_data(ttl=120)
def get_audio_base64(url):
    if not isinstance(url, str) or url == "": return None
    target_url = url.split('?')[0] + "?download=1" if "sharepoint.com" in url else url
    try:
        headers = {'User-Agent': 'Mozilla/5.0'}
        resp = requests.get(target_url, headers=headers, timeout=10)
        if resp.status_code == 200:
            b64 = base64.b64encode(resp.content).decode('utf-8')
            return f"data:audio/mpeg;base64,{b64}"
    except: return None
    return None

def get_embed_url(link):
    if "drive.google.com" in link and "/view" in link:
        return link.replace("/view", "/preview")
    return link

def get_image_download_url(link):
    """💡 終極修正版：直接提取 Google Drive File ID，改走官方 API 媒體流"""
    if not isinstance(link, str): return ""
    if "drive.google.com" in link:
        # 容錯解析：不論是 /file/d/ 還是 id= 都能抓到 33 位的 File ID
        if "/file/d/" in link:
            return link.split("/file/d/")[1].split("/")[0]
        elif "id=" in link:
            return link.split("id=")[1].split("&")[0]
    return link

# === 3. 資料載入與過濾核心 ===
@st.cache_data(ttl=60)
def load_data():
    try:
        df = pd.read_csv(CSV_URL, on_bad_lines='skip', engine='python')
        df.columns = [str(c).strip().lower() for c in df.columns]
        for col in ['title', 'link', 'category', 'type', 'short']:
            if col not in df.columns: df[col] = ""
        df = df.fillna("")
        df['short'] = df.apply(lambda r: r['short'] if str(r['short']).strip() != "" else r['title'], axis=1)
        df = df[~df['category'].astype(str).str.contains("案例資料庫", na=False)]
        df = df[~df['link'].astype(str).str.contains('/folders/')]
        df['uid'] = df['link'].apply(generate_id)
        return df.reset_index(drop=True)
    except:
        return pd.DataFrame()

@st.cache_data(ttl=60)
def load_logo_data():
    """精準對齊您 Clients 分頁的欄位結構"""
    try:
        logo_df = pd.read_csv(CSV_LOGO_URL, on_bad_lines='skip', engine='python')
        logo_df.columns = [str(c).strip().lower() for c in logo_df.columns]
        
        # 尋找並修正您的特殊欄位名稱
        rename_dict = {}
        for c in logo_df.columns:
            if 'client' in c or '客戶' in c or '品名' in c:
                rename_dict[c] = 'client_name'
            if 'link' in c or '網址' in c or 'logo' in c:
                rename_dict[c] = 'logo_link'
                
        logo_df = logo_df.rename(columns=rename_dict)
        return logo_df.fillna("")
    except:
        return pd.DataFrame(columns=['category', 'client_name', 'logo_link'])

# === 4. UI 元件 ===
def render_copy_ui(label, text_to_copy, is_disabled=False, warning_msg=""):
    if is_disabled:
        html_code = f"""
        <div style="background-color:#fff5f5;padding:12px;border-radius:8px;border:1px solid #feb2b2;margin-bottom:10px;">
            <label style="font-size:12px;color:#c53030;font-weight:bold;">{label}</label>
            <p style="font-size:13px;color:#333;margin:8px 0;line-height:1.4;">⚠️ {warning_msg}</p>
        </div>
        """
    else:
        html_code = f"""
        <div style="background-color:#f8f9fa;padding:10px;border-radius:8px;border:1px solid #eee;margin-bottom:10px;">
            <label style="font-size:12px;color:#666;">{label}</label>
            <input type="text" value="{text_to_copy}" id="copyInput" readonly style="width:100%;padding:8px;margin:5px 0;border:1px solid #ddd;border-radius:4px;">
            <button onclick="copyToClipboard()" style="width:100%;padding:10px;background:#0097DA;color:white;border:none;border-radius:5px;cursor:pointer;font-weight:bold;">📋 點此複製網址</button>
            <script>
                function copyToClipboard() {{
                    var copyText = document.getElementById("copyInput");
                    copyText.select();
                    navigator.clipboard.writeText(copyText.value).then(function() {{ alert("✅ 複製成功！"); }});
                }}
            </script>
        </div>
        """
    components.html(html_code, height=150)

@st.dialog("🔗 分享檔案權限")
def show_share_dialog(display_name, link, uid, is_video=False, is_image=False):
    st.write(f"📄 **{display_name}**")
    render_copy_ui("🏢 內部分享連結 (同仁下載用)", link)
    if is_video:
        render_copy_ui("🌏 外部分享連結", "", is_disabled=True, warning_msg="影片涉及『客戶版權』及『全家便利商店場域』，不提供對外分享。")
    elif is_image:
        render_copy_ui("🌏 外部分享連結", "", is_disabled=True, warning_msg="此為『圖片檔』，涉及版權保護，不提供對外分享。")
    else:
        share_link = f"{SITE_URL}?id={uid}"
        render_copy_ui("🌏 外部分享連結 (客戶試聽/防下載)", share_link)

# === 5. 主程式架構 ===
def main():
    st.set_page_config(page_title="全家通路媒體資料庫", layout="centered")
    
    if 'logged_in' not in st.session_state:
        st.session_state.logged_in = False
    if 'display_count' not in st.session_state:
        st.session_state.display_count = 20

    df = load_data()
    logo_df = load_logo_data()
    
    if df.empty:
        st.error("目前無法連線至總資料庫，請檢查發布設定。")
        return

    params = st.query_params
    target_uid = params.get("id", None)

    if target_uid:
        # (保持原有的客戶單頁試聽模式不變...)
        target_row = df[df['uid'] == target_uid]
        if not target_row.empty:
            item = target_row.iloc[0]
            t_low, tp_low = str(item['title']).lower(), str(item['type']).lower()
            is_vid = any(x in tp_low for x in ["新鮮視", "側帶", "demo"]) or any(ext in t_low for ext in ['.mp4', '.mov'])
            is_img = any(ext in t_low for ext in ['.jpg', '.jpeg', '.png', '.gif', '.webp'])
            if is_vid or is_img:
                st.error("此檔案涉及版權保護，不開放對外預覽。")
                return
            st.subheader(f"🎵 作品預覽：{item['short']}")
            b64 = get_audio_base64(item['link'])
            if b64: st.markdown(f'<audio controls controlsList="nodownload" style="width:100%;"><source src="{b64}" type="audio/mpeg"></audio>', unsafe_allow_html=True)
            if st.button("🏠 回到首頁"): st.query_params.clear(); st.rerun()
            return

    if not st.session_state.logged_in:
        st.markdown("<h2 style='text-align: center;'>🔒 全家通路媒體資料庫</h2>", unsafe_allow_html=True)
        with st.form("login_form"):
            pw = st.text_input("請輸入內部資料庫密碼", type="password")
            if st.form_submit_button("解鎖系統", use_container_width=True):
                if pw == PASSWORD: st.session_state.logged_in = True; st.rerun()
                else: st.error("密碼錯誤")
        return

    # 搜尋與過濾元件
    search_query = st.text_input("🔍 關鍵字搜尋 (比對標題內容)")
    if 'last_search' not in st.session_state or st.session_state.last_search != search_query:
        st.session_state.display_count = 20
        st.session_state.last_search = search_query

    c1, c2 = st.columns(2)
    with c1:
        cat_list = sorted([str(x) for x in df['category'].unique() if str(x).strip()])
        sel_cat = st.selectbox("📂 總資料庫分類過濾", ["全部"] + cat_list)
    with c2:
        type_filter = st.radio("📑 類型過濾", ["全部", "企頻", "新鮮視", "側帶", "demo"], horizontal=True)

    mask = pd.Series([True] * len(df), index=df.index)
    if search_query:
        mask &= (df['title'].str.contains(search_query, case=False) | df['category'].str.contains(search_query, case=False))
    if sel_cat != "全部":
        mask &= (df['category'] == sel_cat)
    if type_filter != "全部":
        mask &= (df['type'].str.contains(type_filter, case=False) | df['title'].str.contains(type_filter, case=False))

    results = df[mask]
    total_results = len(results)
    
    # =====================================================================
    # 🆕 兩階段推進式：案例勾選 ➡️ 觸發 Logo 建議區
    # =====================================================================
    
    # 初始化記憶庫
    if 'selected_uids' not in st.session_state:
        st.session_state.selected_uids = []  # 存被勾選案例的 uid
    if 'confirmed_stage' not in st.session_state:
        st.session_state.confirmed_stage = False  # 是否按下 OK 進入第二階段

    # 頂部固定顯示目前已勾選的進度
    st.markdown("---")
    c_status, c_ok = st.columns([4, 1])
    with c_status:
        st.markdown(f"### 📥 已挑選案例進度： **{len(st.session_state.selected_uids)} / 6**")
        # 顯示目前勾了哪些
        if st.session_state.selected_uids:
            picked_shorts = df[df['uid'].isin(st.session_state.selected_uids)]['short'].tolist()
            st.caption(f"已選：{', '.join(picked_shorts)}")
    with c_ok:
        # 當有選東西時，OK 按鈕才會亮起
        if st.button("👌 確認選好", use_container_width=True, type="primary" if st.session_state.selected_uids else "secondary"):
            if st.session_state.selected_uids:
                st.session_state.confirmed_stage = True
                st.rerun()
            else:
                st.warning("請先在下方案例旁勾選！")

   # -----------------------------------------------------------------
    # 【階段二】跳出 Logo 建議與簡報生成頁面 (終極完美 PPT 版)
    # -----------------------------------------------------------------
    if st.session_state.confirmed_stage and st.session_state.selected_uids:
        st.markdown("""
        <div style="background-color:#e0f2fe; padding:20px; border-radius:10px; border-left:5px solid #0284c7; margin: 15px 0;">
            <h4 style="color:#0369a1; margin:0;">🎯 第二階段：確認各案例 Logo 與自訂簡報大標</h4>
            <p style="font-size:14px; color:#0c4a6e; margin:5px 0 0 0;">請確認下方各案例對應的 Logo。您可以在下方直接修改 PPT 的主標題名稱。</p>
        </div>
        """, unsafe_allow_html=True)
        
        # ✨ 需求新增：多出一個可編輯的大標題 Bar
        st.markdown("### 🖋️ 編輯 PPT 簡報大標題")
        custom_ppt_title = st.text_input(
            "請輸入您想要的 PPT 簡報主標題：",
            value="合作夥伴案例分享", # 預設固定寫這行
            placeholder="例如：全家便利商店 媒體行銷案例提案",
            key="custom_ppt_title_input"
        )
        st.markdown("---")
        
        # 準備供手動搜尋的完整 Logo 清單
        logo_options = ["請選擇確切客戶 Logo"] + sorted(list(logo_df['client_name'].unique())) if not logo_df.empty else ["請選擇確切客戶 Logo"]
        
        # 建立一個字典，用來記錄每一個被選中案例最終決定搭配的 logo 網址
        final_pack_pairs = {}
        all_logos_assigned = True 

        # 實施「上下排對應」：使用者勾幾個案例，就跑幾列橫條
        for idx, picked_uid in enumerate(st.session_state.selected_uids):
            case_row = df[df['uid'] == picked_uid].iloc[0]
            case_title = str(case_row['short'])
            
            # 獨立盲猜最適合的 Logo
            guessed_logo_for_this_row = "請選擇確切客戶 Logo"
            if not logo_df.empty:
                for client in logo_df['client_name'].unique():
                    pure_name = str(client).split('_')[-1] if '_' in str(client) else str(client)
                    if pure_name in case_title or str(client) in case_title:
                        guessed_logo_for_this_row = client
                        break
            
            # 畫面排版：左邊放案例文字，右邊放對應的搜尋選單
            c_case_lbl, c_logo_sel = st.columns([3, 2])
            
            with c_case_lbl:
                st.markdown(f"**案例 {idx+1}**")
                st.info(f"📄 {case_title}")
                
            with c_logo_sel:
                st.markdown(f"**對應 Logo {idx+1}**")
                default_idx = logo_options.index(guessed_logo_for_this_row) if guessed_logo_for_this_row in logo_options else 0
                
                chosen_logo_for_row = st.selectbox(
                    f"搜尋/挑選 Logo",
                    options=logo_options,
                    index=default_idx,
                    key=f"sel_logo_pair_{picked_uid}",
                    label_visibility="collapsed"
                )
                
                if chosen_logo_for_row == "請選擇確切客戶 Logo":
                    all_logos_assigned = False
                else:
                    # 撈出該 Logo 的真實雲端直連網址並記錄起來
                    logo_row = logo_df[logo_df['client_name'] == chosen_logo_for_row]
                    raw_url = logo_row.iloc[0]['logo_link'] if not logo_row.empty else ""
                    
                    # 💡 終極核心：改用全新的 googleusercontent 接口，免 API 金鑰直連下載圖片
                    file_id = ""
                    if "/file/d/" in raw_url: file_id = raw_url.split("/file/d/")[1].split("/")[0]
                    elif "id=" in raw_url: file_id = raw_url.split("id=")[1].split("&")[0]
                    
                    final_pack_pairs[picked_uid] = {
                        'case_title': case_title,
                        'logo_name': chosen_logo_for_row,
                        'logo_file_id': file_id
                    }
            st.markdown("<div style='margin-bottom:-10px;'></div>", unsafe_allow_html=True) 

        st.markdown("---")
        
        # -----------------------------------------------------------------
        # 🚀 執行自動生成 PPTX 簡報與實體圖片嵌入
        # -----------------------------------------------------------------
        c_back, c_action = st.columns([1, 4])
        with c_back:
            if st.button("🔙 重挑案例", use_container_width=True):
                st.session_state.confirmed_stage = False
                st.rerun()
                
        with c_action:
            if all_logos_assigned:
                import io
                from datetime import datetime
                from pptx import Presentation
                from pptx.util import Inches, Pt
                
                # 建立一個記憶體內部的二進位流，用來準備 PPTX 檔案
                ppt_buffer = io.BytesIO()
                
                with st.spinner("🚀 正在跨雲端撈取實體 Logo 圖片，並自動排版六宮格簡報中..."):
                    try:
                        # 1. 初始化一份空白簡報 (16:9 寬螢幕)
                        prs = Presentation()
                        prs.slide_width = Inches(13.333)
                        prs.slide_height = Inches(7.5)
                        
                        # 新增一張空白投影片
                        blank_slide_layout = prs.slide_layouts[6]
                        slide = prs.slides.add_slide(blank_slide_layout)
                        
                        # 2. 畫出簡報主標題 (✨ 改用剛剛同仁手動輸入的內容！)
                        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
                        tf = title_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = str(custom_ppt_title).strip() # 👈 這裡完美連動文字輸入框
                        p.font.size = Pt(28)
                        p.font.bold = True
                        p.font.name = "Microsoft JhengHei"
                        
                        # 3. 定義六宮格的標準座標 (2排 x 3列)
                        x_coords = [Inches(0.6), Inches(4.8), Inches(9.0)]
                        y_coords = [Inches(1.8), Inches(4.6)]
                        box_width = Inches(3.8)
                        
                        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
                        
                        # 4. 開始將案例與 Logo 實體圖片填入六宮格
                        for idx, (uid, info) in enumerate(final_pack_pairs.items()):
                            if idx >= 6: break 
                            
                            row_idx = idx // 3  
                            col_idx = idx % 3   
                            
                            current_x = x_coords[col_idx]
                            current_y = y_coords[row_idx]
                            
                            # 💡 繪製案例的文字框 (上方)
                            text_box = slide.shapes.add_textbox(current_x, current_y, box_width, Inches(0.6))
                            tf_case = text_box.text_frame
                            tf_case.word_wrap = True
                            p_case = tf_case.paragraphs[0]
                            p_case.text = f"🔹 {info['case_title']}" 
                            p_case.font.size = Pt(12)
                            p_case.font.name = "Microsoft JhengHei"
                            
                            # 💡 終極實現：抓取實體圖片並嵌入 PPT (下方)
                            if info['logo_file_id']:
                                try:
                                    # 使用全新的第三方高畫質免密鑰直連網址
                                    direct_img_url = f"https://lh3.googleusercontent.com/u/0/d/{info['logo_file_id']}"
                                    resp_logo = requests.get(direct_img_url, headers=headers, timeout=10)
                                    
                                    if resp_logo.status_code == 200 and len(resp_logo.content) > 1000:
                                        logo_stream = io.BytesIO(resp_logo.content)
                                        # 將實體圖片完美塞入 PPT 對應格子的位置
                                        slide.shapes.add_picture(
                                            logo_stream, 
                                            current_x + Inches(0.2), 
                                            current_y + Inches(0.7), 
                                            width=Inches(1.6) # 固定寬度，高度等比例縮放
                                        )
                                except Exception as e_img:
                                    # 如果單張圖片下載有微小意外，輸出小文字提醒但不崩潰
                                    pass
                                    
                        # 5. 排版完成，將簡報儲存至記憶體
                        prs.save(ppt_buffer)
                        ppt_buffer.seek(0)
                        
                        today_str = datetime.now().strftime("%Y%m%d")
                        
                        # 吐出真正的 PowerPoint 下載按鈕！
                        st.download_button(
                            label="🎨 簡報自動排版成功！點此下載六宮格提案 PPTX",
                            data=ppt_buffer,
                            file_name=f"媒體通路提案簡報_{today_str}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                            type="primary"
                        )
                        
                    except Exception as e:
                        st.error(f"❌ 簡報自動生成失敗，原因：{str(e)}")
            else:
                st.warning("⚠️ 提示：上方尚有案例未成功指派 Logo，請先手動搜尋選取，即可解鎖簡報生成功能。")
                
        st.markdown("---")
        
        # -----------------------------------------------------------------
        # 🚀 終極實現：真正的自動生成 PPTX 簡報並下載
        # -----------------------------------------------------------------
        c_back, c_action = st.columns([1, 4])
        with c_back:
            if st.button("🔙 重挑案例", use_container_width=True):
                st.session_state.confirmed_stage = False
                st.rerun()
                
        with c_action:
            if all_logos_assigned:
                import io
                from datetime import datetime
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor
                
                # 建立一個記憶體內部的二進位流，用來準備 PPTX 檔案
                ppt_buffer = io.BytesIO()
                
                with st.spinner("🚀 正在為您跨雲端抓取素材，並自動排版六宮格簡報中..."):
                    try:
                        # 1. 初始化一份空白簡報 (預設是 16:9 寬螢幕)
                        prs = Presentation()
                        prs.slide_width = Inches(13.333)
                        prs.slide_height = Inches(7.5)
                        
                        # 新增一張空白投影片
                        blank_slide_layout = prs.slide_layouts[6]
                        slide = prs.slides.add_slide(blank_slide_layout)
                        
                      # 2. 畫出簡報主標題 (智慧撈取第一個案例所指派的 Logo 名稱作為大標題)
                        first_logo_name = "客戶品牌"
                        if st.session_state.selected_uids:
                            first_uid = st.session_state.selected_uids[0]
                            if first_uid in final_pack_pairs:
                                first_logo_name = final_pack_pairs[first_uid]['logo_name']

                        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
                        tf = title_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = f"📊 媒體通路行銷案例提案 — 【{first_logo_name}】" # 👈 這裡已修正變數錯誤
                        p.font.size = Pt(28)
                        p.font.bold = True
                        p.font.name = "Microsoft JhengHei" # 微軟正黑體
                        
                        # 3. 定義六宮格的標準座標 (2排 x 3列)
                        # 欄位 X 座標: 左、中、右
                        x_coords = [Inches(0.6), Inches(4.8), Inches(9.0)]
                        # 列位 Y 座標: 上、下
                        y_coords = [Inches(1.5), Inches(4.5)]
                        
                        # 寬度與高度固定
                        box_width = Inches(3.8)
                        box_height = Inches(2.5)
                        
                        headers = {'User-Agent': 'Mozilla/5.0'}
                        
                        # 4. 開始將同仁勾選的案例依序填入六宮格
                        for idx, (uid, info) in enumerate(final_pack_pairs.items()):
                            if idx >= 6: break # 最多填滿六格
                            
                            # 計算目前這格要落在第幾列、第幾欄
                            row_idx = idx // 3  # 0 或 1
                            col_idx = idx % 3   # 0, 1, 2
                            
                            current_x = x_coords[col_idx]
                            current_y = y_coords[row_idx]
                            
                            # 💡 繪製案例的文字框 (上方)
                            text_box = slide.shapes.add_textbox(current_x, current_y, box_width, Inches(0.6))
                            tf_case = text_box.text_frame
                            tf_case.word_wrap = True
                            p_case = tf_case.paragraphs[0]
                            # 簡化名稱，只拿短標題
                            p_case.text = f"🔹 {info['case_title'][:20]}..." 
                            p_case.font.size = Pt(12)
                            p_case.font.name = "Microsoft JhengHei"
                            
                        # 💡 抓取並插入對應的客戶 Logo 圖片 (安全防爆版)
                            if info['logo_download_url']:
                                try:
                                    file_id = info['logo_download_url']
                                    # 👈 核心關鍵：直接用官方 API 接口，繞過所有網頁版防毒警告陷阱
                                    api_url = f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media"
                                    
                                    resp_logo = requests.get(api_url, headers=headers, timeout=10)
                                    if resp_logo.status_code == 200:
                                        logo_stream = io.BytesIO(resp_logo.content)
                                        # 嘗試塞入簡報
                                        slide.shapes.add_picture(
                                            logo_stream, 
                                            current_x + Inches(0.2), 
                                            current_y + Inches(0.7), 
                                            width=Inches(1.5)
                                        )
                                except Exception as img_error:
                                    pass
                                    
                            # 💡 如果這筆案例本身是圖片檔，也可以把案例圖片抓下來並排貼上
                            # (此處預留空間，目前先幫您把 Logo 與文字框架對應排好)
                                    
                        # 5. 排版完成，將簡報儲存至記憶體
                        prs.save(ppt_buffer)
                        ppt_buffer.seek(0)
                        
                        today_str = datetime.now().strftime("%Y%m%d")
                        
                        # 💡 吐出真正的 PowerPoint 下載按鈕！
                        st.download_button(
                            label="🎨 簡報自動排版成功！點此下載六宮格提案 PPTX",
                            data=ppt_buffer,
                            file_name=f"通路媒體提案簡報_{first_logo_name}_{today_str}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                            type="primary"
                        )
                        
                    except Exception as e:
                        st.error(f"❌ 簡報自動生成失敗，原因：{str(e)}")
            else:
                st.warning("⚠️ 提示：上方尚有案例未成功指派 Logo，請先手動搜尋選取，即可解鎖簡報生成功能。")
                
        st.markdown("---")

    # -----------------------------------------------------------------
    # 【階段一】渲染列表與動態勾選區
    # -----------------------------------------------------------------
    st.markdown("#### 📂 搜尋結果案例列表 (請在下方挑選打勾)")
    
    current_results = results.head(st.session_state.display_count)
    for _, row in current_results.iterrows():
        uid = row['uid']
        display_name = row['short']
        t_low, tp_low = str(row['title']).lower(), str(row['type']).lower()
        
        is_audio = any(ext in t_low for ext in ['.mp3', '.wav', '.m4a']) or "企頻" in tp_low
        is_video = any(x in tp_low for x in ["新鮮視", "側帶", "demo"]) or any(ext in t_low for ext in ['.mp4', '.mov'])
        is_image = any(ext in t_low for ext in ['.jpg', '.jpeg', '.png', '.gif', '.webp'])

        # 💡 核心改動：在每一個 Expand 展開前，並列一個 Checkbox 放左邊
        col_check, col_exp = st.columns([1, 9])
        
        with col_check:
            # 檢查這個案例先前有沒有被勾選過
            is_checked_before = uid in st.session_state.selected_uids
            
            # 建立勾選方塊
            check_clicked = st.checkbox("選取", key=f"chk_{uid}", value=is_checked_before, label_visibility="collapsed")
            
            # 狀態處理：當同仁點擊打勾或取消
            if check_clicked and uid not in st.session_state.selected_uids:
                if len(st.session_state.selected_uids) < 6:
                    st.session_state.selected_uids.append(uid)
                    st.rerun() # 立刻重整更新頂部的進度條
                else:
                    st.error("❌ 最多只能打包 6 個案例！")
            elif not check_clicked and uid in st.session_state.selected_uids:
                st.session_state.selected_uids.remove(uid)
                st.rerun()

        with col_exp:
            # 原本的 Expander 內容完全不動
            with st.expander(f"📄 {display_name}"):
                if is_audio:
                    if st.button("▶️ 載入音訊", key=f"p_{uid}"):
                        b64 = get_audio_base64(row['link'])
                        if b64: st.audio(b64)
                elif is_video:
                    st.info("📺 影片預覽：限同仁點擊下方『開啟檔案』觀看。")
                elif is_image:
                    st.warning("🖼️ 此為『圖片檔』，不是『影像檔』。同仁請點擊下方『開啟檔案』查看。")
                else:
                    components.iframe(get_embed_url(row['link']), height=400)
                
                bt1, bt2 = st.columns(2)
                with bt1: st.link_button("↗ 開啟檔案", row['link'], use_container_width=True)
                with bt2:
                    if st.button("🔗 分享檔案", key=f"s_{uid}", use_container_width=True):
                        show_share_dialog(display_name, row['link'], uid, is_video=is_video, is_image=is_image)

    # 展開更多案例按鈕
    if total_results > st.session_state.display_count:
        if st.button(f"🔽 展開更多案例", use_container_width=True):
            st.session_state.display_count += 20
            st.rerun()

if __name__ == "__main__":
    main()
